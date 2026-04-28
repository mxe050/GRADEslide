"""
extract.py
使い方: python extract.py input.pptx
出力:
  - extracted_output/extracted.json (テキスト・スタイル・座標・形状情報)
  - extracted_output/images/        (埋め込み画像をオリジナル形式で抽出)
"""

import json
import sys
import traceback
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE


def safe(func, default=None):
    """関数呼び出しを安全に実行(例外を握りつぶす)"""
    try:
        return func()
    except Exception:
        return default


def emu_to_dict(emu_value):
    if emu_value is None:
        return None
    try:
        return {
            "emu": int(emu_value),
            "cm": round(Emu(emu_value).cm, 3),
            "pt": round(Emu(emu_value).pt, 2),
            "inches": round(Emu(emu_value).inches, 3),
        }
    except Exception:
        return None


def color_to_hex(color):
    try:
        if color and color.type is not None:
            return f"#{color.rgb}"
    except Exception:
        pass
    return None


def extract_run(run):
    info = {"text": run.text}
    try:
        font = run.font
        if safe(lambda: font.name):
            info["font_name"] = font.name
        sz = safe(lambda: font.size)
        if sz:
            info["font_size_pt"] = sz.pt
        b = safe(lambda: font.bold)
        if b is not None:
            info["bold"] = b
        it = safe(lambda: font.italic)
        if it is not None:
            info["italic"] = it
        un = safe(lambda: font.underline)
        if un is not None:
            info["underline"] = bool(un)
        hex_color = safe(lambda: color_to_hex(font.color))
        if hex_color:
            info["color"] = hex_color
    except Exception:
        pass
    return info


def extract_paragraph(para):
    return {
        "alignment": safe(lambda: str(para.alignment) if para.alignment else None),
        "level": safe(lambda: para.level, 0),
        "runs": [extract_run(r) for r in para.runs],
    }


def extract_text_frame(tf):
    return {
        "paragraphs": [extract_paragraph(p) for p in tf.paragraphs],
    }


def extract_shape(shape, image_dir, slide_idx, shape_idx):
    info = {
        "shape_id": safe(lambda: shape.shape_id),
        "name": safe(lambda: shape.name),
        "shape_type": safe(lambda: str(shape.shape_type)),
        "left": emu_to_dict(safe(lambda: shape.left)),
        "top": emu_to_dict(safe(lambda: shape.top)),
        "width": emu_to_dict(safe(lambda: shape.width)),
        "height": emu_to_dict(safe(lambda: shape.height)),
    }

    rot = safe(lambda: shape.rotation)
    if rot:
        info["rotation"] = rot

    # テキストフレーム
    if safe(lambda: shape.has_text_frame, False):
        info["text_frame"] = safe(lambda: extract_text_frame(shape.text_frame))

    # 画像
    if safe(lambda: shape.shape_type) == MSO_SHAPE_TYPE.PICTURE:
        try:
            image = shape.image
            ext = image.ext
            filename = f"slide{slide_idx+1}_shape{shape_idx+1}.{ext}"
            filepath = image_dir / filename
            with open(filepath, "wb") as f:
                f.write(image.blob)
            info["image_file"] = str(filepath.as_posix())
            info["image_content_type"] = image.content_type
        except Exception as e:
            info["image_error"] = str(e)

    # 表
    if safe(lambda: shape.has_table, False):
        try:
            table = shape.table
            rows = []
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    cells.append({
                        "text_frame": safe(lambda c=cell: extract_text_frame(c.text_frame))
                    })
                rows.append(cells)
            info["table"] = {
                "rows": len(table.rows),
                "cols": len(table.columns),
                "cells": rows,
            }
        except Exception as e:
            info["table_error"] = str(e)

    # グループ
    if safe(lambda: shape.shape_type) == MSO_SHAPE_TYPE.GROUP:
        try:
            sub_shapes = []
            for j, sub in enumerate(shape.shapes):
                sub_shapes.append(extract_shape(sub, image_dir, slide_idx, shape_idx * 100 + j))
            info["group_shapes"] = sub_shapes
        except Exception as e:
            info["group_error"] = str(e)

    # オートシェイプの種類(プロパティアクセスで例外を出すケースに対応)
    auto_type = safe(lambda: shape.auto_shape_type)
    if auto_type is not None:
        info["auto_shape_type"] = str(auto_type)

    # 塗りつぶし
    fill = safe(lambda: shape.fill)
    if fill is not None:
        ft = safe(lambda: fill.type)
        if ft is not None:
            info["fill_type"] = str(ft)
            hex_color = safe(lambda: color_to_hex(fill.fore_color))
            if hex_color:
                info["fill_color"] = hex_color

    # 枠線
    line = safe(lambda: shape.line)
    if line is not None:
        hex_color = safe(lambda: color_to_hex(line.color))
        if hex_color:
            info["line_color"] = hex_color
        lw = safe(lambda: line.width)
        if lw:
            info["line_width_pt"] = lw.pt

    return info


def extract_slide(slide, image_dir, slide_idx):
    shapes = []
    for j, shape in enumerate(slide.shapes):
        try:
            shapes.append(extract_shape(shape, image_dir, slide_idx, j))
        except Exception as e:
            shapes.append({
                "shape_index": j,
                "extract_error": str(e),
            })

    bg = {}
    try:
        if slide.background and slide.background.fill.type is not None:
            bg["fill_type"] = str(slide.background.fill.type)
            hex_color = safe(lambda: color_to_hex(slide.background.fill.fore_color))
            if hex_color:
                bg["fill_color"] = hex_color
    except Exception:
        pass

    return {
        "slide_index": slide_idx + 1,
        "layout_name": safe(lambda: slide.slide_layout.name),
        "background": bg,
        "shapes": shapes,
    }


def main():
    if len(sys.argv) < 2:
        print("Usage: python extract.py input.pptx")
        sys.exit(1)

    input_path = Path(sys.argv[1])
    if not input_path.exists():
        print(f"File not found: {input_path}")
        sys.exit(1)

    output_dir = Path("extracted_output")
    output_dir.mkdir(exist_ok=True)
    image_dir = output_dir / "images"
    image_dir.mkdir(exist_ok=True)

    prs = Presentation(input_path)

    result = {
        "source_file": str(input_path.name),
        "slide_width": emu_to_dict(prs.slide_width),
        "slide_height": emu_to_dict(prs.slide_height),
        "slide_count": len(prs.slides),
        "slides": [],
    }

    for i, slide in enumerate(prs.slides):
        try:
            result["slides"].append(extract_slide(slide, image_dir, i))
            print(f"  slide {i+1}/{len(prs.slides)} done")
        except Exception as e:
            print(f"  slide {i+1} ERROR: {e}")
            traceback.print_exc()
            result["slides"].append({
                "slide_index": i + 1,
                "extract_error": str(e),
            })

    json_path = output_dir / "extracted.json"
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)

    print(f"OK: JSON saved -> {json_path}")
    print(f"OK: Images saved -> {image_dir}")
    print(f"OK: Total slides -> {result['slide_count']}")


if __name__ == "__main__":
    main()